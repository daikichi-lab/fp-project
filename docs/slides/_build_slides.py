# -*- coding: utf-8 -*-
"""
3日間チャレンジセミナー 台本 → pptx ビルダー
  Day1 / Day2 / Day3 をそれぞれ1ファイルに生成。
  投影スライド = [スライドN] のト書き／スピーカーノート = 読み上げ台詞。
テーマ: 紺×金（高級感・財務/相談役）、明暗サンドイッチ。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ---- palette ----
NAVY   = "0F1A33"   # dark dominant
NAVY2  = "1B2A4A"   # panel on dark
GOLD   = "C6A15B"   # accent
GOLD_D = "A8843F"
GOLD_L = "EAD9B0"
CREAM  = "F6F3EC"   # light content bg
CARD   = "FFFFFF"
INK    = "1E2438"   # body on light
MUTED  = "454B5E"   # 補助テキスト（濃いめ＝可読性重視）
LINE   = "E2DCCF"
WHITE  = "F5F2EA"   # text on dark
GOLDISH= "F4EBD6"   # light gold panel
ACCENT = "A6324A"   # 差し色（バーガンディ）＝同系色を脱する強調用
ACCENT2= "2E6FB0"   # 差し色（ブルー）＝肯定・現金系

SERIF = "Noto Serif CJK JP"   # headers / statements
SANS  = "Noto Sans CJK JP"    # body / tables

EMU_IN = 914400
W, H = 13.333, 7.5

DAY_LABEL = {1: "DAY 1", 2: "DAY 2", 3: "DAY 3"}
DAY_DATE  = {1: "2026.7.14 (火) 20:00–21:00",
             2: "2026.7.16 (木) 20:00–21:00",
             3: "2026.7.18 (土) 20:00–21:00"}

# フッター左のブランド表記（デッキ単位で差し替え可。既定＝3日間チャレンジ）
FOOTER_BRAND = "大吉塾｜FP法人会計マスターコース　無料3日間チャレンジ"


# ---------- low-level helpers ----------
def _font(run, name, size, color, bold=False, italic=False, spacing=None):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = RGBColor.from_string(color)
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", name)
    if spacing is not None:
        rPr.set("spc", str(int(spacing * 100)))


def box(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    return tf


def para(tf, first, align=PP_ALIGN.LEFT, before=0, after=0, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if before: p.space_before = Pt(before)
    if after:  p.space_after = Pt(after)
    if line is not None:
        p.line_spacing = line
    return p


def run(p, text, name, size, color, bold=False, italic=False, spacing=None):
    r = p.add_run(); r.text = text
    _font(r, name, size, color, bold, italic, spacing)
    return r


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE,
         line_color=None, line_w=None, shadow=False, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(color)
    if line_color:
        sp.line.color.rgb = RGBColor.from_string(line_color)
        sp.line.width = Pt(line_w or 1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    if shadow:
        _soft_shadow(sp)
    return sp


def _soft_shadow(sp):
    spPr = sp._element.spPr
    el = spPr.makeelement(qn("a:effectLst"), {})
    sh = el.makeelement(qn("a:outerShdw"),
                        {"blurRad": "90000", "dist": "28000",
                         "dir": "5400000", "rotWithShape": "0"})
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "1A1A1A"})
    alpha = clr.makeelement(qn("a:alpha"), {"val": "18000"})
    clr.append(alpha); sh.append(clr); el.append(sh); spPr.append(el)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text or ""


def footer(slide, idx, total, day, dark=False):
    c = MUTED if not dark else "8A93AD"
    tf = box(slide, 0.9, 7.06, 7.5, 0.3)
    p = para(tf, True)
    run(p, FOOTER_BRAND, SANS, 8.5, c)
    tf2 = box(slide, W - 3.5, 7.06, 2.6, 0.3)
    p2 = para(tf2, True, align=PP_ALIGN.RIGHT)
    run(p2, f"{DAY_LABEL.get(day, '')}  —  {idx:02d} / {total:02d}", SANS, 8.5, c, spacing=0.6)


def motif(slide, kicker, day, dark=False):
    # gold square + kicker label (repeated motif)
    rect(slide, 0.9, 0.66, 0.16, 0.16, GOLD)
    tf = box(slide, 1.2, 0.6, 11.0, 0.34)
    p = para(tf, True)
    run(p, kicker, SANS, 12, GOLD_D if not dark else GOLD_L, bold=True, spacing=1.4)


def base(slide, dark=False):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(NAVY if dark else CREAM)


# ---------- slide builders ----------
def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def title_field(slide, lines, y=1.05, x=0.86, w=12.0, color=INK, size=33, gap=4):
    # auto-shrink long single-line titles so they don't wrap an orphan tail
    longest = max((len(s) for s in lines), default=1)
    fit = int(w * 72 / (1.12 * longest))
    if fit < size:
        size = max(22, fit)
    tf = box(slide, x, y, w, 2.0)
    for i, ln in enumerate(lines):
        p = para(tf, i == 0, after=gap, line=1.04)
        run(p, ln, SERIF, size, color, bold=True)
    return tf


def build_title(prs, d, s, idx, total):
    sl = new_slide(prs); base(sl, dark=True)
    # 装飾：大きな金リング（薄）で奥行き
    ring = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.6), Inches(3.3), Inches(6.4), Inches(6.4))
    ring.fill.background(); ring.line.color.rgb = RGBColor.from_string(NAVY2)
    ring.line.width = Pt(26); ring.shadow.inherit = False
    rect(sl, 0, 0, 0.22, H, GOLD)                       # gold spine
    # eyebrow（金の短いルール＋ラベル）
    rect(sl, 0.97, 1.66, 0.55, 0.045, GOLD)
    tf = box(sl, 0.97, 1.86, 11.0, 0.4)
    run(para(tf, True), s["kicker"], SANS, 12.5, GOLD_L, bold=True, spacing=3.0)
    # title
    tf = box(sl, 0.9, 2.55, 11.5, 2.7)
    for i, ln in enumerate(s["title"]):
        p = para(tf, i == 0, after=8, line=1.1)
        run(p, ln, SERIF, s.get("tsize", 42), WHITE, bold=True)
    if s.get("sub"):
        tf = box(sl, 0.97, 5.3, 10.7, 1.0)
        run(para(tf, True, line=1.35), s["sub"], SANS, 16, GOLD_L)
    rect(sl, 0.97, 6.4, 2.6, 0.025, GOLD)
    tf = box(sl, 0.97, 6.55, 11.0, 0.5)
    run(para(tf, True), s.get("date_line", DAY_DATE.get(d, "") + "　／　ZOOMライブ配信"),
        SANS, 12.5, "C7CEDC")
    footer(sl, idx, total, d, dark=True)
    notes(sl, s.get("notes"))


def build_section(prs, d, s, idx, total):
    sl = new_slide(prs); base(sl, dark=True)
    rect(sl, 0, 0, 0.22, H, GOLD)
    no = s.get("no")
    if no:
        # 背景の大きな章番号（薄）＋前景の金番号
        tf = box(sl, 7.3, 0.8, 5.8, 6.1, anchor=MSO_ANCHOR.MIDDLE)
        run(para(tf, True, align=PP_ALIGN.RIGHT), no, SERIF, 250, NAVY2, bold=True)
        tf = box(sl, 0.97, 1.4, 3.2, 1.3)
        run(para(tf, True), no, SERIF, 78, GOLD, bold=True)
        rect(sl, 1.0, 2.92, 3.6, 0.03, GOLD)
        ky, ty, sy = 3.2, 3.66, 5.55
    else:
        rect(sl, 0.9, 2.35, 0.16, 0.16, GOLD)
        ky, ty, sy = 2.28, 2.85, 5.1
    tf = box(sl, 0.98 if no else 1.2, ky, 11.0, 0.4)
    run(para(tf, True), s.get("kicker", DAY_LABEL.get(d, "")), SANS, 12, GOLD_L, bold=True, spacing=2.5)
    tf = box(sl, 0.9, ty, 11.5, 2.4, anchor=MSO_ANCHOR.TOP)
    for i, ln in enumerate(s["title"]):
        p = para(tf, i == 0, after=6, line=1.1)
        run(p, ln, SERIF, s.get("tsize", 34), WHITE, bold=True)
    if s.get("sub"):
        tf = box(sl, 0.92, sy, 11.2, 1.0)
        run(para(tf, True, line=1.35), s["sub"], SANS, 15, GOLD_L)
    if s.get("avatars"):
        ad = 0.76; ay = 5.55; ax = 8.5
        for lab, col in s["avatars"]:
            _avatar(sl, ax + ad / 2, ay + ad / 2, ad, col)
            tf = box(sl, ax - 0.3, ay + ad + 0.05, ad + 0.6, 0.26)
            run(para(tf, True, align=PP_ALIGN.CENTER), lab, SANS, 9.5, GOLD_L, bold=True)
            ax += ad + 1.1
    footer(sl, idx, total, d, dark=True)
    notes(sl, s.get("notes"))


def build_statement(prs, d, s, idx, total):
    sl = new_slide(prs); base(sl, dark=True)
    rect(sl, 0, 0, 0.22, H, GOLD)
    motif(sl, s.get("kicker", DAY_LABEL[d]), d, dark=True)
    accent = s.get("accent", WHITE)
    tf = box(sl, 1.0, 0, 11.3, H, anchor=MSO_ANCHOR.MIDDLE)
    if s.get("eyebrow"):
        run(para(tf, True, after=14), s["eyebrow"], SANS, 15, GOLD_L)
    lines = s["title"]
    for i, ln in enumerate(lines):
        first = (i == 0) and not s.get("eyebrow")
        p = para(tf, first, after=6, line=1.12)
        run(p, ln, SERIF, s.get("tsize", 38), accent, bold=True)
    if s.get("sub"):
        tf2 = box(sl, 1.02, H - 2.0, 11.0, 1.3)
        run(para(tf2, True, line=1.4), s["sub"], SANS, 15.5, WHITE)
    footer(sl, idx, total, d, dark=True)
    notes(sl, s.get("notes"))


def _lead(slide, text, y=2.0, w=11.5):
    tf = box(slide, 0.9, y, w, 1.0)
    run(para(tf, True, line=1.35), text, SANS, 16, INK)


def build_content(prs, d, s, idx, total):
    sl = new_slide(prs); base(sl)
    motif(sl, s["kicker"], d)
    title_field(sl, s["title"], size=s.get("tsize", 31))
    y = s.get("body_y", 2.25)
    if s.get("lead"):
        _lead(sl, s["lead"], y=y); y += s.get("lead_gap", 0.85)
    items = s.get("items", [])
    n = len(items)
    has_foot = bool(s.get("foot"))
    y_limit = s.get("y_limit", 6.2 if has_foot else 6.5)
    gap = s.get("row_gap", 0.16)
    max_row = s.get("row_h", 1.05)
    row_h = min(max_row, (y_limit - y - gap * (n - 1)) / n) if n else max_row
    for it in items:
        _row(sl, 0.95, y, 11.45, row_h, it)
        y += row_h + gap
    if has_foot:
        tf = box(sl, 1.0, y_limit + 0.1, 11.3, 0.6)
        run(para(tf, True, line=1.3), s["foot"], SANS, 12.5, MUTED, bold=True)
    footer(sl, idx, total, d)
    notes(sl, s.get("notes"))


def _row(slide, x, y, w, h, it):
    # number badge + head + body, on a soft card
    rect(slide, x, y, w, h, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.08, shadow=True)
    rect(slide, x, y, 0.09, h, GOLD)   # left accent strip
    pad = 0.32
    if it.get("n"):
        d = min(h - 0.34, 0.52)
        rect(slide, x + pad, y + (h - d) / 2, d, d, NAVY, shape=MSO_SHAPE.OVAL)
        tf = box(slide, x + pad, y + (h - d) / 2 - 0.02, d, d, anchor=MSO_ANCHOR.MIDDLE)
        run(para(tf, True, align=PP_ALIGN.CENTER), str(it["n"]), SERIF, 20, GOLD, bold=True)
        tx = x + pad + d + 0.3
    elif it.get("icon"):
        d = min(h - 0.34, 0.52)
        rect(slide, x + pad, y + (h - d) / 2, d, d, GOLDISH, shape=MSO_SHAPE.OVAL)
        tfc = box(slide, x + pad, y + (h - d) / 2 - 0.02, d, d, anchor=MSO_ANCHOR.MIDDLE)
        run(para(tfc, True, align=PP_ALIGN.CENTER), it["icon"], SERIF, 20, GOLD_D, bold=True)
        tx = x + pad + d + 0.3
    else:
        tx = x + pad
    tw = x + w - tx - 0.3
    tf = box(slide, tx, y, tw, h, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, True, line=1.18, after=2)
    run(p, it["head"], SANS, it.get("hsize", 17), INK, bold=True)
    if it.get("body"):
        p2 = para(tf, False, line=1.18)
        run(p2, it["body"], SANS, it.get("bsize", 13), MUTED, bold=True)


def build_cards(prs, d, s, idx, total):
    sl = new_slide(prs); base(sl)
    motif(sl, s["kicker"], d)
    title_field(sl, s["title"], size=s.get("tsize", 31))
    y0 = s.get("body_y", 2.45)
    if s.get("lead"):
        _lead(sl, s["lead"], y=y0); y0 += 0.9
    items = s["items"]; n = len(items)
    gx = 0.95; gw = 11.45
    if n == 4:
        cols, rows = 2, 2
    elif n <= 3:
        cols, rows = n, 1
    else:
        cols, rows = 3, (n + 2) // 3
    gap = 0.3
    cw = (gw - gap * (cols - 1)) / cols
    has_foot = bool(s.get("foot"))
    bottom = 6.35 if has_foot else 6.75
    avail_h = bottom - y0
    ch = (avail_h - gap * (rows - 1)) / rows
    ch = min(ch, 2.1)   # avoid over-tall cards w/ empty bottom
    dh = 15 if cols >= 3 else 17               # narrower cards → smaller heading
    db = 12.5 if cols >= 3 else 13
    for it in items:
        it.setdefault("hsize", dh); it.setdefault("bsize", db)
    for i, it in enumerate(items):
        r, c = divmod(i, cols)
        x = gx + c * (cw + gap); y = y0 + r * (ch + gap)
        _card(sl, x, y, cw, ch, it)
    if has_foot:
        tf = box(sl, 1.0, bottom + 0.16, 11.3, 0.6)
        run(para(tf, True, line=1.25), s["foot"], SANS, 12, MUTED, bold=True)
    footer(sl, idx, total, d)
    notes(sl, s.get("notes"))


def _card(slide, x, y, w, h, it):
    rect(slide, x, y, w, h, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.06, shadow=True)
    rect(slide, x, y, w, 0.1, GOLD)
    pad = 0.3
    cy = y + 0.34
    if it.get("icon") and not it.get("n"):
        isz = 0.52
        rect(slide, x + pad, y + 0.3, isz, isz, GOLDISH, shape=MSO_SHAPE.OVAL)
        tfi = box(slide, x + pad, y + 0.28, isz, isz, anchor=MSO_ANCHOR.MIDDLE)
        run(para(tfi, True, align=PP_ALIGN.CENTER), it["icon"], SERIF, 22, GOLD_D, bold=True)
        tfh = box(slide, x + pad + isz + 0.16, y + 0.3, w - pad * 2 - isz - 0.16, isz,
                  anchor=MSO_ANCHOR.MIDDLE)
        run(para(tfh, True, line=1.08), it["head"], SANS, it.get("hsize", 16), INK, bold=True)
        if it.get("body"):
            tfb = box(slide, x + pad, y + 0.3 + isz + 0.12, w - pad * 2,
                      h - (0.3 + isz + 0.12) - 0.18)
            run(para(tfb, True, line=1.3), it["body"], SANS, it.get("bsize", 12.5), MUTED, bold=True)
        return
    if it.get("n"):
        rect(slide, x + pad, cy, 0.5, 0.5, GOLDISH, shape=MSO_SHAPE.OVAL)
        tf = box(slide, x + pad, cy - 0.02, 0.5, 0.5, anchor=MSO_ANCHOR.MIDDLE)
        run(para(tf, True, align=PP_ALIGN.CENTER), str(it["n"]), SERIF, 19, GOLD_D, bold=True)
        cy += 0.66
    tf = box(slide, x + pad, cy, w - pad * 2, h - (cy - y) - 0.25)
    p = para(tf, True, line=1.2, after=5)
    run(p, it["head"], SANS, it.get("hsize", 17), INK, bold=True)
    if it.get("body"):
        run(para(tf, False, line=1.32), it["body"], SANS, it.get("bsize", 13), MUTED, bold=True)


def build_stat(prs, d, s, idx, total):
    sl = new_slide(prs); base(sl)
    motif(sl, s["kicker"], d)
    if s.get("icon"):
        rect(sl, 0.9, 1.04, 0.6, 0.6, GOLDISH, shape=MSO_SHAPE.OVAL, shadow=True)
        tfi = box(sl, 0.9, 1.02, 0.6, 0.6, anchor=MSO_ANCHOR.MIDDLE)
        run(para(tfi, True, align=PP_ALIGN.CENTER), s["icon"], SERIF, 24, GOLD_D, bold=True)
        title_field(sl, s["title"], x=1.72, w=10.4, size=s.get("tsize", 31))
    else:
        title_field(sl, s["title"], size=s.get("tsize", 31))
    top = s.get("body_y", 2.4)
    bottom = 6.5
    bandH = bottom - top
    # ---- 左：論点を白カードに接地（—羅列をやめ、金ノードで） ----
    lx, lw = 0.9, 6.4
    rect(sl, lx, top, lw, bandH, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05, shadow=True)
    rect(sl, lx, top, 0.1, bandH, GOLD)
    if s.get("lead"):
        run(para(box(sl, lx + 0.5, top + 0.18, lw - 0.8, 0.5), True, line=1.25),
            s["lead"], SANS, 13.5, GOLD_D, bold=True)
    pts = s.get("points", [])
    n = max(len(pts), 1)
    pad_top = 0.55 if s.get("lead") else 0.3
    inner = bandH - pad_top - 0.25
    rowh = min(1.05, inner / n)
    gy = top + pad_top + (inner - rowh * n) / 2
    for it in pts:
        cyr = gy + rowh / 2
        rect(sl, lx + 0.5, cyr - 0.1, 0.2, 0.2, GOLD, shape=MSO_SHAPE.OVAL)
        tf = box(sl, lx + 0.95, gy, lw - 1.25, rowh, anchor=MSO_ANCHOR.MIDDLE)
        run(para(tf, True, line=1.2), it["head"], SANS, it.get("hsize", 15.5), INK, bold=True)
        if it.get("body"):
            run(para(tf, False, line=1.2), it["body"], SANS, 12.5, MUTED, bold=True)
        gy += rowh
    # ---- 右：数字パネル（主役）＋差し色タグ ----
    px, pw = 7.7, 4.7
    rect(sl, px, top, pw, bandH, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05, shadow=True)
    rect(sl, px, top, pw, 0.14, GOLD)
    cx = px + pw / 2
    if s.get("tag"):
        ttext, tcol = s["tag"]
        tw = 0.32 * len(ttext) + 0.6
        rect(sl, cx - tw / 2, top + 0.42, tw, 0.46, tcol,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        tf = box(sl, cx - tw / 2, top + 0.40, tw, 0.46, anchor=MSO_ANCHOR.MIDDLE)
        run(para(tf, True, align=PP_ALIGN.CENTER), ttext, SANS, 12.5, "FFFFFF", bold=True)

    def _lines(box_y, box_h, val, size, color, bold=False, lh=1.2):
        tf = box(sl, px + 0.25, box_y, pw - 0.5, box_h, anchor=MSO_ANCHOR.MIDDLE)
        lst = [val] if isinstance(val, str) else val
        for i, ln in enumerate(lst):
            run(para(tf, i == 0, align=PP_ALIGN.CENTER, line=lh), ln, SANS, size, color, bold=bold)

    if s.get("stat_label"):
        _lines(top + 1.0, 0.6, s["stat_label"], 13.5, GOLD_L, lh=1.15)
    tf = box(sl, px + 0.15, top + 1.5, pw - 0.3, 1.35, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, True, align=PP_ALIGN.CENTER, line=1.0), s["stat"], SERIF,
        s.get("stat_size", 54), GOLD, bold=True)
    if s.get("stat_sub"):
        _lines(top + 2.95, 0.85, s["stat_sub"], 13.5, WHITE, bold=True, lh=1.3)
    if s.get("stat_note"):
        tf = box(sl, px + 0.2, top + bandH - 0.42, pw - 0.4, 0.3)
        run(para(tf, True, align=PP_ALIGN.CENTER), s["stat_note"], SANS, 10, "8A93AD")
    footer(sl, idx, total, d)
    notes(sl, s.get("notes"))


def _style_table(tbl):
    # strip default banding
    tbl.first_row = False; tbl.horz_banding = False


def _cell(cell, text, font=SANS, size=12.5, color=INK, bold=False,
          fill=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    cell.vertical_anchor = anchor
    cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.12)
    cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
    if fill is None:
        cell.fill.background()
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor.from_string(fill)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    _font(r, font, size, color, bold)


def build_table_kv(prs, d, s, idx, total):
    """2列の見出し+値テーブル（区分表/コース表など）。rows: [(left,right,highlight?)]"""
    sl = new_slide(prs); base(sl)
    motif(sl, s["kicker"], d)
    title_field(sl, s["title"], size=s.get("tsize", 31))
    y0 = s.get("body_y", 2.35)
    if s.get("lead"):
        _lead(sl, s["lead"], y=y0); y0 += 0.85
    head = s["head"]; rows = s["rows"]
    nrow = len(rows) + 1
    tx, tw = 1.0, 11.3
    th = min(0.6 * nrow, 6.4 - y0)
    gtbl = sl.shapes.add_table(nrow, 2, Inches(tx), Inches(y0),
                               Inches(tw), Inches(th)).table
    _style_table(gtbl)
    gtbl.columns[0].width = Inches(tw * s.get("c0", 0.5))
    gtbl.columns[1].width = Inches(tw * (1 - s.get("c0", 0.5)))
    cmp = s.get("compare")   # 右列を“魅力side”としてゴールドのレーンに
    _cell(gtbl.cell(0, 0), head[0], SANS, 13.5, WHITE, bold=True, fill=NAVY)
    _cell(gtbl.cell(0, 1), head[1], SANS, 13.5,
          NAVY if cmp else WHITE, bold=True, fill=(GOLD_D if cmp else NAVY))
    for i, rw in enumerate(rows, start=1):
        hl = len(rw) > 2 and rw[2]
        base_fill = CARD if i % 2 else CREAM
        fill0 = GOLDISH if hl else base_fill
        _cell(gtbl.cell(i, 0), rw[0], SANS, s.get("fs", 12.5),
              NAVY if hl else INK, bold=hl, fill=fill0)
        if cmp:
            _cell(gtbl.cell(i, 1), rw[1], SANS, s.get("fs", 12.5),
                  NAVY, bold=True, fill=(GOLD_L if hl else "FAF3E0"),
                  align=s.get("c1align", PP_ALIGN.LEFT))
        else:
            _cell(gtbl.cell(i, 1), rw[1], SANS, s.get("fs", 12.5),
                  INK, bold=hl, fill=fill0, align=s.get("c1align", PP_ALIGN.LEFT))
    if s.get("foot"):
        tf = box(sl, 1.0, y0 + th + 0.18, 11.3, 0.6)
        run(para(tf, True, line=1.3), s["foot"], SANS, 12.5, MUTED, bold=True)
    footer(sl, idx, total, d)
    notes(sl, s.get("notes"))


def build_bs(prs, d, s, idx, total):
    """貸借対照表（T字）。"""
    sl = new_slide(prs); base(sl)
    motif(sl, s["kicker"], d)
    title_field(sl, s["title"], size=s.get("tsize", 30))
    tf = box(sl, 0.95, 1.9, 11.4, 0.45)
    run(para(tf, True), s.get("lead", ""), SANS, 14, MUTED)
    y0 = 2.3
    rows = s["rows"]  # list of (la, lv, ra, rv, lhl, rhl)
    nrow = len(rows) + 1
    tx, tw = 1.0, 11.3
    bottom = 6.35 if s.get("foot") else 6.6
    th = min(0.42 * nrow, bottom - y0)
    tbl = sl.shapes.add_table(nrow, 4, Inches(tx), Inches(y0),
                              Inches(tw), Inches(th)).table
    _style_table(tbl)
    for c, frac in zip(range(4), (0.34, 0.16, 0.34, 0.16)):
        tbl.columns[c].width = Inches(tw * frac)
    hd = ["資産の部", "百万円", "負債・純資産の部", "百万円"]
    for c in range(4):
        _cell(tbl.cell(0, c), hd[c], SANS, 12.5, WHITE, bold=True, fill=NAVY,
              align=PP_ALIGN.RIGHT if c in (1, 3) else PP_ALIGN.LEFT)
    for i, rw in enumerate(rows, start=1):
        la, lv, ra, rv = rw[0], rw[1], rw[2], rw[3]
        lhl = len(rw) > 4 and rw[4]; rhl = len(rw) > 5 and rw[5]
        base_fill = CARD if i % 2 else CREAM
        for c, (txt, hl) in enumerate([(la, lhl), (lv, lhl), (ra, rhl), (rv, rhl)]):
            fill = GOLDISH if hl else base_fill
            sub = txt.startswith("（")
            _cell(tbl.cell(i, c), txt, SANS, 12 if not sub else 11,
                  NAVY if hl else (MUTED if sub else INK),
                  bold=hl, fill=fill,
                  align=PP_ALIGN.RIGHT if c in (1, 3) else PP_ALIGN.LEFT)
    if s.get("foot"):
        tf = box(sl, 1.0, y0 + th + 0.16, 11.3, 0.5)
        run(para(tf, True, line=1.3), s["foot"], SANS, 12.5, MUTED, bold=True)
    footer(sl, idx, total, d)
    notes(sl, s.get("notes"))


def build_price(prs, d, s, idx, total):
    sl = new_slide(prs); base(sl, dark=True)
    rect(sl, 0, 0, 0.22, H, GOLD)
    motif(sl, s["kicker"], d, dark=True)
    tf = box(sl, 1.0, 1.5, 11.0, 0.6)
    run(para(tf, True), s["title"][0], SERIF, 26, WHITE, bold=True)
    # big price
    rect(sl, 3.4, 2.7, 6.5, 2.4, NAVY2, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.06, shadow=True)
    rect(sl, 3.4, 2.7, 6.5, 0.12, GOLD)
    tf = box(sl, 3.4, 2.7, 6.5, 2.4, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, True, align=PP_ALIGN.CENTER, after=4), s["plabel"], SANS, 16, GOLD_L)
    run(para(tf, False, align=PP_ALIGN.CENTER), s["price"], SERIF, 64, GOLD, bold=True)
    tf = box(sl, 1.0, 5.4, 11.3, 1.2, anchor=MSO_ANCHOR.TOP)
    run(para(tf, True, align=PP_ALIGN.CENTER, line=1.4), s["sub"], SANS, 15, WHITE)
    footer(sl, idx, total, d, dark=True)
    notes(sl, s.get("notes"))


# ---- お金の地図（ブロックパズル型・未来会計の中核ビジュアル）----
# (fill_on, text_on, fill_dim, text_dim)
MM_COLORS = {
    "売上":   ("1B2A4A", "FFFFFF", "CDD4E1", "3D4A66"),
    "変動費": ("9AA0AE", "FFFFFF", "DCDDE3", "5A6173"),
    "粗利":   ("C6A15B", "1B2A4A", "EAD9B0", "6B5A2E"),
    "固定費": ("33405C", "FFFFFF", "CDD4E1", "3D4A66"),
    "残り":   ("C6A15B", "1B2A4A", "EAD9B0", "6B5A2E"),
}


def _mm_seg(sl, x, y, w, h, name, label, on):
    fill_on, txt_on, fill_dim, txt_dim = MM_COLORS[name]
    fill = fill_on if on else fill_dim
    rect(sl, x, y, w, h, fill, shadow=(on and h >= 0.5))
    gold = name in ("粗利", "残り")
    if h < 0.46:
        # block too thin → label outside (to the right), readable on cream
        tf = box(sl, x + w + 0.12, y - 0.16, 2.4, h + 0.32, anchor=MSO_ANCHOR.MIDDLE)
        p = para(tf, True, line=1.0)
        run(p, name, SANS, 13, (GOLD_D if gold else NAVY), bold=True)
        if label:
            run(p, "  " + label, SANS, 12, MUTED)
        return
    txt = txt_on if on else txt_dim
    tf = box(sl, x + 0.08, y, w - 0.16, h, anchor=MSO_ANCHOR.MIDDLE)
    if h < 0.72:
        p = para(tf, True, align=PP_ALIGN.CENTER, line=1.0)
        run(p, name, SANS, 13 if on else 12, txt, bold=on)
        if label:
            run(p, "  " + label, SANS, 11.5, txt)
    else:
        p = para(tf, True, align=PP_ALIGN.CENTER, line=1.05)
        run(p, name, SANS, 16 if on else 13, txt, bold=on)
        if label:
            run(para(tf, False, align=PP_ALIGN.CENTER), label, SANS, 12.5, txt)


def build_moneymap(prs, d, s, idx, total):
    """お金の地図：売上→変動費(素通り)→粗利→固定費→残り の積み上げ図。
       fields: amounts(dict or None), highlight(売上/変動費/粗利/固定費/残り/all/None),
               show_values(bool), unit(str), lead, notes"""
    sl = new_slide(prs); base(sl)
    motif(sl, s["kicker"], d)
    title_field(sl, s["title"], size=s.get("tsize", 29))
    if s.get("lead"):
        _lead(sl, s["lead"], y=2.0)
    a = s.get("amounts") or {"売上": 100, "変動費": 55, "粗利": 45, "固定費": 34, "残り": 11}
    hi = s.get("highlight")
    sv = s.get("show_values", False)
    unit = s.get("unit", "")
    top, H = 2.75, 3.3
    scale = H / a["売上"]
    w, gap = 2.3, 0.55
    x1 = 2.4; x2 = x1 + w + gap; x3 = x2 + w + gap
    on = lambda nm: hi in (None, "all") or hi == nm
    lab = lambda nm: (f'{a[nm]}{unit}' if sv else "")
    h_vc = a["変動費"] * scale; h_gp = a["粗利"] * scale
    h_fc = a["固定費"] * scale; h_nk = a["残り"] * scale
    _mm_seg(sl, x1, top, w, H, "売上", lab("売上"), on("売上"))
    _mm_seg(sl, x2, top, w, h_vc, "変動費", lab("変動費"), on("変動費"))
    _mm_seg(sl, x2, top + h_vc, w, h_gp, "粗利", lab("粗利"), on("粗利"))
    _mm_seg(sl, x3, top + h_vc, w, h_fc, "固定費", lab("固定費"), on("固定費"))
    _mm_seg(sl, x3, top + h_vc + h_fc, w, h_nk, "残り", lab("残り"), on("残り"))
    # captions
    tf = box(sl, x2 - 0.3, top - 0.36, w + 0.6, 0.3)
    run(para(tf, True, align=PP_ALIGN.CENTER), "↑ ヨソに素通り（仕入・材料）", SANS, 11, "5A6173")
    tf = box(sl, x3, top + h_vc - 0.34, w, 0.3)
    run(para(tf, True, align=PP_ALIGN.CENTER), "↓ 粗利の中身", SANS, 11, "5A6173")
    tf = box(sl, x3 - 0.3, top + H + 0.2, w + 0.6, 0.32)
    run(para(tf, True, align=PP_ALIGN.CENTER), "ここが“残る／残らない”", SANS, 11.5, GOLD_D, bold=True)
    footer(sl, idx, total, d)
    notes(sl, s.get("notes"))


# ---- 三表連動図（國貞A型を自前再現：PL→BS / CF→BS現金）----
L3_ORANGE = "D58A36"   # 利益系
L3_BLUE = "2E6FB0"     # 現金系
L3_GRY = "AAB0BE"      # その他（沈める）
L3_GOLD = "C6A15B"


def _l3seg(sl, x, y, w, h, color, label, txt="FFFFFF", size=10.5, bold=False):
    rect(sl, x, y, w, h, color)
    if label:
        tf = box(sl, x + 0.03, y, w - 0.06, h, anchor=MSO_ANCHOR.MIDDLE)
        run(para(tf, True, align=PP_ALIGN.CENTER), label, SANS, size, txt, bold=bold)


def _l3cap(sl, x, y, w, text, color, align=PP_ALIGN.CENTER, size=11, bold=True):
    tf = box(sl, x, y, w, 0.3)
    run(para(tf, True, align=align), text, SANS, size, color, bold=bold)


def _l3arrow(sl, x1, y1, x2, y2, color):
    cxn = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                  Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cxn.line.color.rgb = RGBColor.from_string(color)
    cxn.line.width = Pt(2.75)
    ln = cxn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    cxn.shadow.inherit = False
    return cxn


def build_link3(prs, d, s, idx, total):
    """三表連動図：PL当期純利益→BS利益剰余金（橙）／CF→BS現金（青）。"""
    sl = new_slide(prs); base(sl)
    motif(sl, s["kicker"], d)
    title_field(sl, s["title"], size=s.get("tsize", 29))
    top = 2.75
    # ---- PL（左） ----
    px, pw = 1.0, 2.1
    _l3cap(sl, px - 0.1, top - 0.62, pw + 0.2, "損益計算書（PL）", INK)
    _l3cap(sl, px - 0.1, top - 0.34, pw + 0.2, "売上高", MUTED, size=10, bold=False)
    pl = [("売上総利益", 1.45, L3_GRY, "FFFFFF"), ("営業利益", 0.5, L3_GRY, "FFFFFF"),
          ("経常利益", 0.45, L3_GRY, "FFFFFF"), ("当期純利益", 0.55, L3_ORANGE, "FFFFFF")]
    y = top
    for nm, h, c, t in pl:
        _l3seg(sl, px, y, pw, h, c, nm, t, bold=(c == L3_ORANGE)); y += h
    pl_net_cy = y - 0.275          # 当期純利益の中心
    # ---- CF（中） ----
    cx, cw = 3.85, 1.7
    _l3cap(sl, cx - 0.15, top - 0.62, cw + 0.3, "キャッシュフロー計算書", INK, size=10)
    cf = [("営業CF ③", 0.8, L3_GOLD), ("投資CF ②", 0.7, "D9BE86"), ("財務CF ①", 0.7, L3_GOLD)]
    y = top + 0.25
    for nm, h, c in cf:
        _l3seg(sl, cx, y, cw, h, c, nm, "1B2A4A", 10, True); y += h
    cf_out_cy = top + 0.25 + 0.4   # 上端付近（現金へ向かう）
    # ---- BS（右・2列） ----
    bx, colw, gp = 7.5, 2.15, 0.18
    bh_top = top
    _l3cap(sl, bx, top - 0.62, colw, "資産", INK)
    _l3cap(sl, bx + colw + gp, top - 0.62, colw, "負債・純資産", INK)
    _l3cap(sl, bx - 0.2, top - 0.34, colw * 2 + gp + 0.4, "貸借対照表（BS）", MUTED, size=10, bold=False)
    # 資産列
    _l3seg(sl, bx, bh_top, colw, 0.7, L3_BLUE, "現金", "FFFFFF", 11, True)
    _l3seg(sl, bx, bh_top + 0.7, colw, 0.9, L3_GRY, "流動資産")
    _l3seg(sl, bx, bh_top + 1.6, colw, 1.35, "B9BFC9", "固定資産")
    # 負債・純資産列
    bx2 = bx + colw + gp
    _l3seg(sl, bx2, bh_top, colw, 0.85, L3_GRY, "流動負債")
    _l3seg(sl, bx2, bh_top + 0.85, colw, 1.0, "B9BFC9", "固定負債")
    _l3seg(sl, bx2, bh_top + 1.85, colw, 0.55, L3_GRY, "純資産")
    _l3seg(sl, bx2, bh_top + 2.4, colw, 0.55, L3_ORANGE, "利益剰余金", "FFFFFF", 10, True)
    res_cy = bh_top + 2.675
    cash_cy = bh_top + 0.35
    # ---- 矢印＋凡例（段階表示：stage 0=箱のみ / 1=橙 / 2=橙青 / 3=完成） ----
    st = s.get("stage", 3)
    if st >= 1:
        _l3arrow(sl, px + pw + 0.02, pl_net_cy, bx2 - 0.02, res_cy, L3_ORANGE)   # 純利益→剰余金
        _l3cap(sl, 1.0, 6.30, 11.3,
               "① 当期純利益 → 利益剰余金（儲けは“純資産”に積み上がる）", L3_ORANGE, PP_ALIGN.LEFT, 12.5)
    if st >= 2:
        _l3arrow(sl, cx + cw + 0.02, cf_out_cy, bx - 0.02, cash_cy, L3_BLUE)     # CF→現金
        _l3cap(sl, 1.0, 6.62, 11.3,
               "② キャッシュフローの現金 → BSの“現金”（残った現金は一致する）", L3_BLUE, PP_ALIGN.LEFT, 12.5)
    footer(sl, idx, total, d)
    notes(sl, s.get("notes"))


# ---- チャット（人型アイコン＋吹き出し）----
def _avatar(sl, cx, cy, d, badge):
    """丸バッジ＋白い人型シルエット（頭＋肩）。"""
    rect(sl, cx - d / 2, cy - d / 2, d, d, badge, shape=MSO_SHAPE.OVAL, shadow=True)
    hd = d * 0.30
    rect(sl, cx - hd / 2, cy - d * 0.27, hd, hd, CARD, shape=MSO_SHAPE.OVAL)
    bw, bh = d * 0.58, d * 0.42
    rect(sl, cx - bw / 2, cy + d * 0.03, bw, bh, CARD, shape=MSO_SHAPE.OVAL)


def _bubble(sl, x, y, w, h, fill, side, text, sub=None, size=14, tcolor=INK):
    """吹き出し：角丸＋三角テール（avatar 方向を指す）。sub＝本音などの第2行。"""
    rect(sl, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.14, shadow=True)
    ts = 0.2
    ty = y + h * 0.5 - ts / 2
    if side == "left":
        tri = sl.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  Inches(x - ts * 0.5), Inches(ty), Inches(ts), Inches(ts))
        tri.rotation = 270
    else:
        tri = sl.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  Inches(x + w - ts * 0.5), Inches(ty), Inches(ts), Inches(ts))
        tri.rotation = 90
    tri.fill.solid(); tri.fill.fore_color.rgb = RGBColor.from_string(fill)
    tri.line.fill.background(); tri.shadow.inherit = False
    tf = box(sl, x + 0.3, y, w - 0.6, h, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, True, line=1.18), text, SANS, size, tcolor, bold=True)
    if sub:
        run(para(tf, False, line=1.22, before=4), sub, SANS, max(size - 3, 11.5), MUTED, bold=True)


def build_chat(prs, d, s, idx, total):
    sl = new_slide(prs); base(sl)
    motif(sl, s["kicker"], d)
    title_field(sl, s["title"], size=s.get("tsize", 30))
    msgs = s["items"]; n = len(msgs)
    y = s.get("body_y", 2.2)
    has_foot = bool(s.get("foot"))
    bottom = 6.2 if has_foot else 6.55
    gap = s.get("row_gap", 0.18)
    h = min(s.get("bubble_h", 1.15), (bottom - y - gap * (n - 1)) / n)
    L, R = 1.0, 12.33
    bw = s.get("bubble_w", 7.7)
    ad = min(0.54, h * 0.52)
    for m in msgs:
        spk = m.get("speaker", "")
        side = m.get("side") or ("right" if spk.startswith("社長") else "left")
        badge = m.get("color") or (GOLD if spk.startswith("社長") else NAVY2)
        fill = m.get("fill") or (GOLDISH if side == "right" else CARD)
        cy = y + h * 0.42
        if side == "left":
            acx = L + ad / 2; bx = L + ad + 0.34
            nx, nw, nal = bx + bw + 0.3, R - (bx + bw + 0.3), PP_ALIGN.LEFT
        else:
            acx = R - ad / 2; bx = R - ad - 0.34 - bw
            nx, nw, nal = L + 0.1, bx - 0.34 - (L + 0.1), PP_ALIGN.RIGHT
        _avatar(sl, acx, cy, ad, badge)
        lab = box(sl, acx - 0.55, cy + ad / 2 + 0.04, 1.1, 0.22)
        run(para(lab, True, align=PP_ALIGN.CENTER), spk, SANS, 8.5, MUTED, bold=True)
        _bubble(sl, bx, y, bw, h, fill, side, m["text"], sub=m.get("sub"), size=m.get("size", 14))
        if m.get("note") and nw > 1.0:
            nb = box(sl, nx, y, nw, h, anchor=MSO_ANCHOR.MIDDLE)
            run(para(nb, True, align=nal, line=1.25), m["note"], SANS, 12, MUTED, bold=True)
        y += h + gap
    if has_foot:
        tf = box(sl, 1.0, bottom + 0.12, 11.3, 0.5)
        run(para(tf, True, line=1.3), s["foot"], SANS, 12.5, MUTED, bold=True)
    footer(sl, idx, total, d)
    notes(sl, s.get("notes"))


# ---- 比較マトリクス（項目×選択肢／○△×評価＋説明）----
def _mark(sl, cx, cy, size, kind):
    if kind == "good":          # 青の○リング
        o = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - size / 2), Inches(cy - size / 2),
                                Inches(size), Inches(size))
        o.fill.background(); o.line.color.rgb = RGBColor.from_string(ACCENT2)
        o.line.width = Pt(3.6); o.shadow.inherit = False
    elif kind == "warn":        # グレーの△
        t = sl.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(cx - size / 2),
                                Inches(cy - size / 2), Inches(size), Inches(size))
        t.fill.solid(); t.fill.fore_color.rgb = RGBColor.from_string("AFB5C0")
        t.line.fill.background(); t.shadow.inherit = False
    else:                       # 赤の×（2本のバー）
        for ang in (45, -45):
            r = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx - size / 2),
                                    Inches(cy - 0.035), Inches(size), Inches(0.075))
            r.rotation = ang
            r.fill.solid(); r.fill.fore_color.rgb = RGBColor.from_string(ACCENT)
            r.line.fill.background(); r.shadow.inherit = False


def _mcell(sl, x, y, w, h, text, fill, color, bold=False, align=PP_ALIGN.CENTER, size=13):
    rect(sl, x, y, w, h, fill, line_color=LINE, line_w=0.75)
    tf = box(sl, x + 0.16, y, w - 0.32, h, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tf, True, align=align, line=1.12), text, SANS, size, color, bold=bold)


def build_matrix(prs, d, s, idx, total):
    sl = new_slide(prs); base(sl)
    motif(sl, s["kicker"], d)
    title_field(sl, s["title"], size=s.get("tsize", 31))
    opts, rows = s["options"], s["rows"]
    n = len(opts)
    x0, TW = 0.9, 11.5
    cw = s.get("crit_w", 2.3)
    ow = (TW - cw) / n
    y0 = s.get("body_y", 2.25)
    hh = 0.64
    rh = (6.45 - y0 - hh) / len(rows)
    # header
    _mcell(sl, x0, y0, cw, hh, s.get("crit_head", "比較項目"), "EFEADC", MUTED, bold=True, size=12.5)
    for j, opt in enumerate(opts):
        ox = x0 + cw + j * ow
        good = opt.get("good")
        _mcell(sl, ox, y0, ow, hh, opt["name"], "F2E7CB" if good else "ECEDF1",
               GOLD_D if good else NAVY, bold=True, size=14)
    # body
    for i, rw in enumerate(rows):
        ry = y0 + hh + i * rh
        _mcell(sl, x0, ry, cw, rh, rw["label"], "EFEADC", INK, bold=True, size=13.5)
        for j, (mk, text) in enumerate(rw["cells"]):
            ox = x0 + cw + j * ow
            cfill = "FCF8EF" if opts[j].get("good") else CARD
            rect(sl, ox, ry, ow, rh, cfill, line_color=LINE, line_w=0.75)
            _mark(sl, ox + ow / 2, ry + rh * 0.33, 0.34, mk)
            tf = box(sl, ox + 0.2, ry + rh * 0.54, ow - 0.4, rh * 0.44, anchor=MSO_ANCHOR.TOP)
            run(para(tf, True, align=PP_ALIGN.CENTER, line=1.12), text, SANS, 12.5, INK, bold=True)
    if s.get("foot"):
        tf = box(sl, 1.0, y0 + hh + rh * len(rows) + 0.16, 11.3, 0.5)
        run(para(tf, True, line=1.3), s["foot"], SANS, 12.5, MUTED, bold=True)
    footer(sl, idx, total, d)
    notes(sl, s.get("notes"))


# ---- アジェンダ（章の地図）----
def build_agenda(prs, d, s, idx, total):
    sl = new_slide(prs); base(sl)
    motif(sl, s.get("kicker", "CONTENTS"), d)
    title_field(sl, s["title"], size=s.get("tsize", 31))
    items = s["items"]; n = len(items)
    y = s.get("body_y", 2.55)
    rh = min(1.25, (6.4 - y) / n)
    for k, it in enumerate(items):
        ry = y + k * rh
        tf = box(sl, 1.0, ry, 1.5, rh, anchor=MSO_ANCHOR.MIDDLE)
        run(para(tf, True), it["no"], SERIF, 38, GOLD, bold=True)
        tf = box(sl, 2.55, ry, 9.7, rh, anchor=MSO_ANCHOR.MIDDLE)
        run(para(tf, True, line=1.12, after=2), it["head"], SANS, 18, INK, bold=True)
        if it.get("body"):
            run(para(tf, False, line=1.2), it["body"], SANS, 12.5, MUTED, bold=True)
        if k < n - 1:
            rect(sl, 1.0, ry + rh, 11.3, 0.012, LINE)
    footer(sl, idx, total, d)
    notes(sl, s.get("notes"))


# ---- タイムライン／プロセス（横並びの番号ノード）----
def build_timeline(prs, d, s, idx, total):
    sl = new_slide(prs); base(sl)
    motif(sl, s["kicker"], d)
    title_field(sl, s["title"], size=s.get("tsize", 31))
    if s.get("lead"):
        _lead(sl, s["lead"], y=2.15)
    steps = s["items"]; n = len(steps)
    ny = s.get("node_y", 3.5)
    x0, x1 = 1.2, 12.1
    span = x1 - x0
    gap = span / n
    rect(sl, x0 + gap * 0.5, ny - 0.02, span - gap, 0.045, GOLD)   # 接続線
    dd = 0.86
    for i, it in enumerate(steps):
        cx = x0 + gap * (i + 0.5)
        rect(sl, cx - dd / 2, ny - dd / 2, dd, dd, NAVY, shape=MSO_SHAPE.OVAL, shadow=True)
        rect(sl, cx - dd / 2, ny - dd / 2, dd, dd, NAVY, shape=MSO_SHAPE.OVAL)
        tf = box(sl, cx - dd / 2, ny - dd / 2 - 0.02, dd, dd, anchor=MSO_ANCHOR.MIDDLE)
        run(para(tf, True, align=PP_ALIGN.CENTER), str(i + 1), SERIF, 28, GOLD, bold=True)
        tf = box(sl, cx - gap / 2 + 0.1, ny + 0.7, gap - 0.2, 0.55, anchor=MSO_ANCHOR.TOP)
        run(para(tf, True, align=PP_ALIGN.CENTER, line=1.12), it["head"], SANS, 15.5, INK, bold=True)
        if it.get("body"):
            tf = box(sl, cx - gap / 2 + 0.12, ny + 1.32, gap - 0.24, 1.4, anchor=MSO_ANCHOR.TOP)
            run(para(tf, True, align=PP_ALIGN.CENTER, line=1.22), it["body"], SANS, 11.5, MUTED, bold=True)
    if s.get("foot"):
        tf = box(sl, 1.0, 6.5, 11.3, 0.4)
        run(para(tf, True, align=PP_ALIGN.CENTER), s["foot"], SANS, 12.5, MUTED, bold=True)
    footer(sl, idx, total, d)
    notes(sl, s.get("notes"))


BUILDERS = {
    "title": build_title, "section": build_section, "statement": build_statement,
    "content": build_content, "cards": build_cards, "stat": build_stat,
    "table_kv": build_table_kv, "bs": build_bs, "price": build_price,
    "moneymap": build_moneymap, "link3": build_link3, "chat": build_chat,
    "matrix": build_matrix, "agenda": build_agenda, "timeline": build_timeline,
}


def build_day(day, slides, fname, subtitle):
    prs = Presentation()
    prs.slide_width = Emu(int(W * EMU_IN))
    prs.slide_height = Emu(int(H * EMU_IN))
    prs.core_properties.title = subtitle
    prs.core_properties.author = "大吉塾 / 藤山泰成"
    total = len(slides)
    for i, s in enumerate(slides, start=1):
        BUILDERS[s["kind"]](prs, day, s, i, total)
    path = os.path.join(OUT_DIR, fname)
    prs.save(path)
    print("saved", path, total, "slides")


# ============================================================
#  DATA は別モジュールから
# ============================================================
if __name__ == "__main__":
    from _slides_data import DAY1, DAY2, DAY3
    build_day(1, DAY1, "Day1_黒字なのに金が残らない.pptx", "Day1 黒字なのに、金が残らない")
    build_day(2, DAY2, "Day2_決算書はこう読む.pptx", "Day2 決算書は、こう読む ― 4か所を読んで分析する")
    build_day(3, DAY3, "Day3_初回面談から契約まで_会いに行ける.pptx", "Day3 初回面談から“契約”まで ― この流れで、会いに行ける")
